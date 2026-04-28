"""
build-pitch-deck.py

Builds the 7-slide investor pitch deck for Lenga Maps.

Output: pitch/lenga-maps-pitch.pptx

Brand:
  - background: dark navy (#0D2B45)
  - accent:     gold (#C9A227)
  - text:       white / soft white

Slide canvas: 13.333 x 7.5 in (16:9 widescreen).

Re-runnable. Idempotent. No commentary inside slides — the deck is delivered
to a live audience and shouldn't read like an AI-generated treatment.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT     = Path(__file__).resolve().parent.parent
PITCH    = ROOT / "pitch"
LOGO     = ROOT / "public" / "images" / "branding" / "logo.png"
PHOTO    = PITCH / "founder.jpg"
OUT_PPTX = PITCH / "lenga-maps-pitch.pptx"

# ── Brand palette ────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x2B, 0x45)
NAVY_DEEP   = RGBColor(0x0A, 0x1F, 0x33)  # for slight panel contrast on s2
GOLD        = RGBColor(0xC9, 0xA2, 0x27)
GOLD_SOFT   = RGBColor(0xDA, 0xB5, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_SOFT  = RGBColor(0xE6, 0xED, 0xF5)  # body text
WHITE_MUTE  = RGBColor(0xA8, 0xB8, 0xC8)  # captions / tertiary
PILL_BG     = RGBColor(0x18, 0x36, 0x52)  # tech-stack chip background
PILL_BORDER = RGBColor(0x2A, 0x4B, 0x6C)

# ── Canvas ───────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────

def add_solid_bg(slide, color: RGBColor) -> None:
    """Paint the entire slide with a solid colour."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(
    slide,
    left, top, width, height,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = WHITE,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font_name: str = "Calibri",
    line_spacing: float = 1.15,
):
    """Place a text box with one paragraph."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(
    slide,
    left, top, width, height,
    lines,                                # list[(text, kwargs)]
    *,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font_name: str = "Calibri",
    line_spacing: float = 1.2,
    paragraph_space_before: float = 4,
):
    """Place a text box with multiple paragraphs, each with its own style."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    for i, item in enumerate(lines):
        text = item["text"]
        size = item.get("size", 18)
        bold = item.get("bold", False)
        color = item.get("color", WHITE)
        space_before = item.get("space_before", paragraph_space_before)
        font = item.get("font", font_name)
        ls = item.get("line_spacing", line_spacing)

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get("align", align)
        p.line_spacing = ls
        if i > 0:
            p.space_before = Pt(space_before)

        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_pill(slide, left, top, text: str, *, padding_x=0.18, height=0.42, font_size=13):
    """A rounded chip used in the tech-stack columns."""
    # rough text-width estimate; PowerPoint will autosize-fit in practice
    char_w = font_size * 0.0085  # inches per character (rough)
    width_in = max(0.9, len(text) * char_w + padding_x * 2)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top,
        Inches(width_in), Inches(height),
    )
    shape.adjustments[0] = 0.5  # full pill curvature
    shape.fill.solid()
    shape.fill.fore_color.rgb = PILL_BG
    shape.line.color.rgb = PILL_BORDER
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = WHITE_SOFT
    return shape, width_in


def flow_pills(slide, items, *, left_in: float, top_in: float, max_width_in: float,
               row_gap=0.18, pill_gap=0.14, font_size=13, height=0.42):
    """Lay pills in flow rows; wrap when row width exceeds max_width_in."""
    cur_x = left_in
    cur_y = top_in
    for label in items:
        # estimate width identically to add_pill
        char_w = font_size * 0.0085
        w = max(0.9, len(label) * char_w + 0.18 * 2)
        if cur_x - left_in + w > max_width_in and cur_x > left_in:
            cur_x = left_in
            cur_y += height + row_gap
        add_pill(slide, Inches(cur_x), Inches(cur_y), label,
                 height=height, font_size=font_size)
        cur_x += w + pill_gap
    return cur_y + height  # bottom-y consumed (in inches)


def add_accent_dot(slide, left_in: float, top_in: float, size_in: float = 0.16):
    """Small gold square dot used as a bullet marker."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left_in), Inches(top_in),
        Inches(size_in), Inches(size_in),
    )
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = GOLD
    s.shadow.inherit = False
    return s


def add_image(slide, path, left, top, width=None, height=None):
    return slide.shapes.add_picture(str(path), left, top,
                                    width=width, height=height)


# ──────────────────────────────────────────────────────────────────────────────
# Slide builders
# ──────────────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s, NAVY)

    # Logo, large, upper third, centred
    logo_w = Inches(2.2)
    logo_top = Inches(1.55)
    logo_left = Inches((13.333 - 2.2) / 2)
    add_image(s, LOGO, logo_left, logo_top, width=logo_w)

    # "Lenga Maps" headline — generous box height so descenders don't collide
    # with the tagline below.
    add_text(s, Inches(0), Inches(3.85), Inches(13.333), Inches(1.4),
             "Lenga Maps",
             font_size=64, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font_name="Georgia",
             line_spacing=1.05)

    # tagline in gold — pushed down to clear the headline cap-line
    add_text(s, Inches(0), Inches(5.30), Inches(13.333), Inches(0.6),
             "Africa's GIS Data Backbone",
             font_size=26, bold=False, color=GOLD,
             align=PP_ALIGN.CENTER, font_name="Calibri")

    # date
    add_text(s, Inches(0), Inches(6.00), Inches(13.333), Inches(0.4),
             "Investor Pitch  ·  April 2026",
             font_size=15, color=WHITE_MUTE,
             align=PP_ALIGN.CENTER)

    # Founder credit, bottom-right
    add_text(s, Inches(7.0), Inches(7.0), Inches(6.0), Inches(0.35),
             "Mulenga Chilufya  ·  Founder & Lead Engineer",
             font_size=11, color=WHITE_MUTE,
             align=PP_ALIGN.RIGHT)


def slide_founder(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s, NAVY)

    # Photo on left, ~5 in wide
    photo_w = Inches(4.6)
    photo_h = Inches(5.6)
    photo_left = Inches(0.7)
    photo_top  = Inches(0.95)
    add_image(s, PHOTO, photo_left, photo_top, width=photo_w, height=photo_h)

    # Name + subtitle (right column)
    rx = 5.7  # x-start of right column (in inches)
    rw = 7.0  # right-column width

    add_text(s, Inches(rx), Inches(0.85), Inches(rw), Inches(0.85),
             "Mulenga Chilufya",
             font_size=42, bold=True, color=WHITE, font_name="Georgia")

    add_text(s, Inches(rx), Inches(1.65), Inches(rw), Inches(0.45),
             "Founder & Lead Engineer  ·  Environmental Engineer",
             font_size=14, color=GOLD)

    # ROLES section header
    add_text(s, Inches(rx), Inches(2.30), Inches(rw), Inches(0.35),
             "ROLES",
             font_size=11, bold=True, color=WHITE_MUTE)

    role_lines = [
        ("Founder & Lead Engineer", "Lenga Maps", "Nov 2024 – present"),
        ("Global Chapters Director", "ThinkOcean Society", "Aug 2025 – present"),
        ("Co-Founder & COO", "Clinzed Ltd", "Sep 2023 – Dec 2025"),
    ]
    role_top = 2.70
    for i, (title, org, dates) in enumerate(role_lines):
        y = role_top + i * 0.62
        add_accent_dot(s, rx, y + 0.13)
        add_text(s, Inches(rx + 0.32), Inches(y), Inches(rw - 0.32), Inches(0.32),
                 f"{title}  —  {org}",
                 font_size=16, bold=True, color=WHITE)
        add_text(s, Inches(rx + 0.32), Inches(y + 0.32), Inches(rw - 0.32), Inches(0.28),
                 dates,
                 font_size=12, color=WHITE_MUTE)

    # EDUCATION section header
    edu_top = role_top + len(role_lines) * 0.62 + 0.10
    add_text(s, Inches(rx), Inches(edu_top), Inches(rw), Inches(0.35),
             "EDUCATION",
             font_size=11, bold=True, color=WHITE_MUTE)

    edu_lines = [
        ("B.Eng. Environmental Engineering", "The Copperbelt University"),
        ("[High school — to fill in]", ""),
    ]
    for i, (deg, inst) in enumerate(edu_lines):
        y = edu_top + 0.40 + i * 0.42
        add_accent_dot(s, rx, y + 0.10)
        line = deg if not inst else f"{deg}  —  {inst}"
        add_text(s, Inches(rx + 0.32), Inches(y), Inches(rw - 0.32), Inches(0.32),
                 line,
                 font_size=14, color=WHITE_SOFT)

    # Affiliations strip (bottom)
    aff_top = 6.85
    add_text(s, Inches(0.7), Inches(aff_top), Inches(11.95), Inches(0.45),
             "Affiliations:  Middlesex Univ. London Accelerator  ·  Commonwealth Engineer's Council  ·  "
             "SADC Groundwater Management Institute  ·  Engineering Institution of Zambia  ·  "
             "African Climate Creatives Fellowship",
             font_size=10, color=WHITE_MUTE)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s, NAVY)

    # Eyebrow label
    add_text(s, Inches(0.9), Inches(0.85), Inches(8), Inches(0.35),
             "THE PROBLEM",
             font_size=12, bold=True, color=GOLD)

    # Pull-quote, large, ~80% width
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(2.6),
             "“According to the World Bank, no region on Earth has a "
             "data-infrastructure index as low as Africa’s.”",
             font_size=40, bold=True, color=WHITE,
             font_name="Georgia", line_spacing=1.18)

    # Sub-statement
    add_text(s, Inches(0.9), Inches(5.20), Inches(11.5), Inches(1.8),
             "Every climate, agricultural and infrastructure decision on the continent is "
             "being made on data that doesn’t exist, isn’t accessible, or isn’t trusted.",
             font_size=22, color=WHITE_SOFT,
             line_spacing=1.32)


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s, NAVY)

    # eyebrow
    add_text(s, Inches(0), Inches(1.05), Inches(13.333), Inches(0.4),
             "OUR SOLUTION",
             font_size=12, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)

    # giant centred sentence
    add_text(s, Inches(1.2), Inches(2.45), Inches(10.95), Inches(3.3),
             "We’re building Africa’s largest and most centralised "
             "environmental GIS database.",
             font_size=46, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER,
             font_name="Georgia",
             line_spacing=1.22)

    # logo at bottom
    logo_w = Inches(0.9)
    add_image(s, LOGO, Inches((13.333 - 0.9) / 2), Inches(6.30), width=logo_w)


def slide_traction(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s, NAVY)

    add_text(s, Inches(0.9), Inches(0.85), Inches(8), Inches(0.4),
             "TRACTION",
             font_size=12, bold=True, color=GOLD)

    add_text(s, Inches(0.9), Inches(1.30), Inches(11.5), Inches(0.85),
             "What’s already shipping",
             font_size=38, bold=True, color=WHITE, font_name="Georgia")

    bullets = [
        ("Platform live in production",
         "lengamaps.com — public, browsable, downloadable."),
        ("11+ datasets shipped",
         "Covering all 54 African countries — climate, hydrology, soil, land cover and more."),
        ("Revenue model running",
         "Mobile-money payments verified end-to-end; recurring billing wired in."),
        ("System architecture built ground-up",
         "No off-the-shelf GIS platform underneath — every layer engineered in-house."),
        ("Middlesex University London accelerator",
         "Accepted onto a competitive programme that pitches founders to UK investors."),
    ]
    top = 2.55
    line_h = 0.85
    for i, (head, sub) in enumerate(bullets):
        y = top + i * line_h
        add_accent_dot(s, 1.0, y + 0.18, size_in=0.22)
        add_text(s, Inches(1.55), Inches(y), Inches(11), Inches(0.42),
                 head,
                 font_size=20, bold=True, color=WHITE)
        add_text(s, Inches(1.55), Inches(y + 0.40), Inches(11), Inches(0.40),
                 sub,
                 font_size=14, color=WHITE_MUTE)


def slide_tech(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s, NAVY)

    add_text(s, Inches(0.9), Inches(0.85), Inches(8), Inches(0.4),
             "TECH STACK",
             font_size=12, bold=True, color=GOLD)
    add_text(s, Inches(0.9), Inches(1.30), Inches(11.5), Inches(0.85),
             "Built with",
             font_size=38, bold=True, color=WHITE, font_name="Georgia")

    # Two columns. Left: GIS. Right: PLATFORM. Top of pills aligned.
    col_top = 2.85
    col_w = 5.7

    # GIS column
    add_text(s, Inches(0.9), Inches(col_top), Inches(col_w), Inches(0.4),
             "GIS",
             font_size=15, bold=True, color=GOLD)
    gis_items = [
        "QGIS", "ArcGIS", "GDAL / OGR", "PostGIS",
        "Global Mapper", "Google Earth Engine", "Blender",
        "Python · rasterio", "Python · geopandas",
    ]
    flow_pills(s, gis_items,
               left_in=0.9, top_in=col_top + 0.55,
               max_width_in=col_w + 0.3,
               font_size=13, height=0.46)

    # PLATFORM column
    add_text(s, Inches(7.0), Inches(col_top), Inches(col_w), Inches(0.4),
             "PLATFORM",
             font_size=15, bold=True, color=GOLD)
    plat_items = [
        "Next.js 14", "TypeScript", "Supabase",
        "Cloudflare R2", "Vercel", "TailwindCSS",
        "Agentic AI orchestration · Anthropic Claude",
    ]
    flow_pills(s, plat_items,
               left_in=7.0, top_in=col_top + 0.55,
               max_width_in=col_w + 0.3,
               font_size=13, height=0.46)


def slide_close(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s, NAVY)

    # logo top-centre
    logo_w = Inches(1.6)
    add_image(s, LOGO, Inches((13.333 - 1.6) / 2), Inches(0.95), width=logo_w)

    # giant closer — generous box height; this line wraps to two lines on
    # 16:9 at 50pt so we need ~1.9in to seat them without overflowing.
    add_text(s, Inches(0.5), Inches(2.65), Inches(12.33), Inches(2.1),
             "Let’s build Africa’s data backbone.",
             font_size=50, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font_name="Georgia",
             line_spacing=1.18)

    # contact lines — pushed below the closer's actual rendered footprint
    add_multiline(s,
        Inches(0), Inches(5.05), Inches(13.333), Inches(1.85),
        [
            {"text": "cmulenga672@gmail.com",  "size": 22, "bold": True, "color": GOLD,        "align": PP_ALIGN.CENTER},
            {"text": "+260 965 699 359",       "size": 18, "color": WHITE_SOFT,               "align": PP_ALIGN.CENTER, "space_before": 8},
            {"text": "lengamaps.com",          "size": 18, "color": WHITE_SOFT,               "align": PP_ALIGN.CENTER, "space_before": 4},
        ],
    )

    # footer
    add_text(s, Inches(0), Inches(7.0), Inches(13.333), Inches(0.35),
             "Mulenga Chilufya  ·  Founder, Lenga Maps",
             font_size=11, color=WHITE_MUTE, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# Driver
# ──────────────────────────────────────────────────────────────────────────────

def main() -> int:
    if not LOGO.exists():
        print(f"missing logo: {LOGO}", file=sys.stderr); return 1
    if not PHOTO.exists():
        print(f"missing founder photo: {PHOTO}", file=sys.stderr); return 1

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_founder(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_traction(prs)
    slide_tech(prs)
    slide_close(prs)

    PITCH.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"wrote {OUT_PPTX}  ({OUT_PPTX.stat().st_size // 1024} KB)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
